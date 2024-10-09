import PyPDF2
from PIL import Image
import speech_recognition as sr 
import moviepy.editor as mp
import math
import librosa
import soundfile
from moviepy.editor import VideoFileClip
from docx import Document as DocxDocument
from PIL import Image
from moviepy.editor import VideoFileClip
import mutagen 
import os
import pandas as pd
from pptx import Presentation
from PIL import Image
import pytesseract
from django.core.files.storage import default_storage


import warnings
warnings.filterwarnings('ignore')
from collections import Counter
import numpy as np
import re
import string
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.decomposition import TruncatedSVD
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.stem import SnowballStemmer
import nltk
nltk.download('punkt_tab')



def extract_text_from_pdf(pdf_file):
    text = ""
    try:
        reader = PyPDF2.PdfReader(pdf_file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        text = "Error extracting text from PDF"
    return text


def extract_text_from_word( word_file):
    text = []
    try:
        doc = DocxDocument(word_file)
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        extracted_text = "\n".join(text)
    except Exception as e:
        print(f"Error extracting text from Word: {e}")
        extracted_text = "Error extracting text from Word"
    return extracted_text

def extract_text_from_text( text_file):
        text = ""
        try:
            text_file.seek(0)  # Ensure the file pointer is at the beginning
            text = text_file.read().decode('utf-8')  # Decode the bytes to a string
        except Exception as e:
            print(f"Error extracting text from text file: {e}")
            text = "Error extracting text from text file"
        return text

def extract_text_from_audio( audio_file):

    #     # Convert m4a to wav
    # if audio_file.name.endswith('m4a'):
    #     audio_segment = AudioSegment.from_file(audio_file, format='m4a')
    #     wav_file_path = 'output/audio.wav'
    #     audio_segment.export(wav_file_path, format='wav')
    #     audio_file = wav_file_path  # Update the audio_file to the new wav file

    audio_samples, sample_rate = librosa.load(audio_file, sr=None)
    audio_duration_sec = librosa.get_duration(y=audio_samples, sr=sample_rate)

    # Create a directory to store subsegment WAV files
    output_dir = 'output/subsegments'
    os.makedirs(output_dir, exist_ok=True)

    subsegment_paths = []
    if audio_duration_sec > 180:
        num_subsegments = math.ceil(audio_duration_sec / 180)
        samples_per_subsegment = len(audio_samples) // num_subsegments
        
        for i in range(num_subsegments):
            start_sample = i * samples_per_subsegment
            end_sample = (i + 1) * samples_per_subsegment
            subsegment_samples = audio_samples[start_sample:end_sample]
            subsegment_path = os.path.join(output_dir, f'subsegment_{i + 1}.wav')
            soundfile.write(subsegment_path, subsegment_samples, sample_rate)
            subsegment_paths.append(subsegment_path)
    else:
        audio_path = os.path.join('output', 'audio.wav')
        soundfile.write(audio_path, audio_samples, sample_rate)
        subsegment_paths = [audio_path]

    r = sr.Recognizer()
    results = []
    for audio in subsegment_paths:
        with sr.AudioFile(audio) as source:
            audio_data = r.record(source)
            try:
                # result = r.recognize_google(audio_data, language='en-US')
                result = r.recognize_google(audio_data, language='ar-AR')
                results.append(result)
            except sr.UnknownValueError:
                print("Google Speech Recognition could not understand audio")
            except sr.RequestError as e:
                print(f"Could not request results from Google Speech Recognition service; {e}")

    return " ".join(results)  # Combine results into a single string

def extract_text_from_video( video_file):
    text = ""

    if hasattr(video_file, 'path'):
        video_file_path = video_file.path

    # Convert video to audio
    clip = mp.VideoFileClip(video_file_path)
    clip.audio.write_audiofile(r"output\converted.wav")
    audio_file_path = os.path.join('output', "converted.wav")

    text = extract_text_from_audio(audio_file_path)
    return text

def extract_text_from_excel( excel_file):
    text = ""
    try:
        # Excel
        if excel_file.name.endswith(('xlsx')):
            df = pd.read_excel(excel_file, sheet_name=None) 
            for sheet_name, sheet_data in df.items():
                text += f"Sheet: {sheet_name}\n"
                text += sheet_data.to_string(index=False) + "\n\n"

        # CSV
        elif excel_file.name.endswith(('csv')):
            df = pd.read_csv(excel_file)  
            text += df.to_string(index=False) + "\n\n"

        
        else:
            text = "Unsupported file type."

    except Exception as e:
        print(f"Error extracting text from file: {e}")
        text = "Error extracting text from file"
    
    return text

def extract_text_from_powerpoint( powerpoint_file):
    """Extract text from a PowerPoint file."""
    text = ""
    try:
        # open file PowerPoint
        presentation = Presentation(powerpoint_file)

        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n" 

    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        text = "Error extracting text from PowerPoint"

    return text

def extract_text_from_image( image_file):
    """Extract text from an image file."""
    text = ""
    try:
        image = Image.open(image_file)
        text = pytesseract.image_to_string(image, lang='ara') 

    except Exception as e:
        print(f"Error extracting text from image: {e}")
        text = "Error extracting text from image"

    return text



def details_image(image_file):
    """Get details of an image file."""
    # Open the image file
    try:
        image = Image.open(image_file)

        # Get image details
        width, height = image.size
        format = image.format
        mode = image.mode

        # Create a details
        details = (
            f"Image Details:\n"
            f"- Format: {format}\n"
            f"- Mode: {mode}\n"
            f"- Dimensions: {width}x{height}px"
        )
    except Exception as e:
        print(f"Error getting details from file: {e}")
        details = "Error getting details from file"
    return details

def details_video(video_file):
    """Get details of a video file."""
    # Open the video file
    if hasattr(video_file, 'path'):
        video_file_path = video_file.path
    try:
        clip = VideoFileClip(video_file_path)

        # Get video details
        duration = clip.duration  # Duration in seconds
        fps = clip.fps  # Frames per second
        width, height = clip.size  # Dimensions
        num_frames = int(clip.fps * clip.duration)  # Total number of frames

        # Create a details
        details = (
            f"Video Details:\n"
            f"- Duration: {duration:.2f} seconds\n"
            f"- Frames per Second (FPS): {fps}\n"
            f"- Dimensions: {width}x{height}px\n"
            f"- Total Frames: {num_frames}"
        )
    except Exception as e:
        print(f"Error getting details from file: {e}")
        details = "Error getting details from file"

    return details

def details_audio(audio_file):
    """Get details of an audio file."""
    # Open the audio file
    try:
        audio = mutagen.File(audio_file)
        # Get audio details
        duration = audio.info.length  # Duration in seconds
        bitrate = audio.info.bitrate  # Bitrate in bps
        channels = audio.info.channels  # Number of channels

        # Create a details
        details = (
            f"Audio Details:\n"
            f"- Duration: {duration:.2f} seconds\n"
            f"- Bitrate: {bitrate // 1000} kbps\n"
            f"- Channels: {channels}"
        )
    except Exception as e:
        print(f"Error getting details from file: {e}")
        details = "Error getting details from file"
    return details

def details_document( doc_file):
    """Get details of a document file."""
    # Get the size of the document directly from the FieldFile
    try:
        size = doc_file.size  # Size in bytes
        file_name = os.path.basename(doc_file.name)  # Use .name to get the file name

        # Create a details string
        details = (
            f"Document Details:\n"
            f"- File Name: {file_name}\n"
            f"- Size: {size / (1024 * 1024):.2f} MB\n"  # Convert to MB
        )
        
    except Exception as e:
        print(f"Error getting details from file: {e}")
        details = "Error getting details from file"
    return details

def details_excel( excel_file):
    """Get details of a file based on its type."""
    details = ""
    try:
        #   Excel
        if excel_file.name.endswith(('xlsx')):
            df = pd.read_excel(excel_file, sheet_name=None) 
            num_sheets = len(df)
            details = (
                f"Excel Details:\n"
                f"- Number of Sheets: {num_sheets}\n"
            )

        #   CSV
        elif excel_file.name.endswith(('csv')):
            dff = pd.read_csv(excel_file)  
            num_rows, num_cols = dff.shape
            num_sheets = len(dff)
            details = (
                f"CSV Details:\n"
                f"- Number of Rows: {num_rows}\n"
                f"- Number of Columns: {num_cols}\n"
                f"- Number of Sheets: {num_sheets}\n"
            )
        else:
            details = "Unsupported file type."

    except Exception as e:
        print(f"Error getting details from file: {e}")
        details = "Error getting details from file"
    
    return details



def description_pdf(self , pdf_file):
    self
def description_Word(self, doc_file):
    self
def description_txt(self, txt_file):
    self
def description_audio(self, audio_file):
    self
def description_video(self,video_file):
    self
def description_image(self, image_file):
    self
def description_excel(self, excel_file):
    self 
def description_powerpoint(self, powerpoint_file):
    self


def clean_text( text):
    """
    Preprocesses text by removing punctuation, converting to lowercase,
    stemming, and removing stopwords.
    """
    text = text.lower()
    words = word_tokenize(text)
    words = [word for word in words if word not in string.punctuation]
    stemmer = SnowballStemmer('english')
    words = [stemmer.stem(word) for word in words]
    stop_words = set(stopwords.words('english'))
    words = [word for word in words if word not in stop_words]
    return ' '.join(words)

def summarize_text(self, text, max_words=150):
    """
    Summarizes a text document using Latent Semantic Analysis (LSA).

    Args:
        text (str): The text to be summarized.
        max_words (int, optional): The maximum number of words in the summary.
            Defaults to 150.

    Returns:
        str: The summarized text.
    """

    cleaned_text = clean_text(text)
    corpus = sent_tokenize(cleaned_text)

    vectorizer = TfidfVectorizer(stop_words='english')
    tdm = vectorizer.fit_transform(corpus)

    decomposer = TruncatedSVD(n_components=3, random_state=0)
    reduced = decomposer.fit_transform(tdm)

    max_cols = reduced.argmax(axis=1)
    included = set()
    summarized_sentences = []
    summary = []

    already_included = set()
    total_words = 0
    selected = []

    for i in max_cols:
        if i not in already_included:
            already_included.add(i)
            selected.append((i, corpus[i]))  # Store sentence index and text
            total_words += len(corpus[i].split(' '))
            if total_words > max_words:
                break

    summary = [s[1] for s in sorted(selected, key=lambda x: x[0])]  # Sort by sentence index
    return ' '.join(summary)

# # Example usage
# text = read_text(path)  # Assuming your read_text function is defined elsewhere
# summary = summarize_text(text)
# print(summary)
